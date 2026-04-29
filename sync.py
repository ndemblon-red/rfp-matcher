"""Sync: read PPTX slides between section dividers, infer metadata locally, upsert to database."""
import hashlib
import logging
import os
import re
import shutil
import time

from pptx import Presentation
from dotenv import load_dotenv

from db import upsert_case_study, log_sync_run, get_case_study_count, get_case_study_by_hash, update_slide_num

load_dotenv()
logger = logging.getLogger(__name__)


# ── Content hashing ───────────────────────────────────────────────────────────

def _hash_content(text):
    """Return a 12-hex-character SHA-256 digest of the slide text for change detection."""
    return hashlib.sha256((text or "").encode("utf-8")).hexdigest()[:12]


# ── Slide patterns ─────────────────────────────────────────────────────────────

# Matches the "CASE STUDIES | Project Name" heading line on content slides
_HEADING_RE      = re.compile(r'CASE\s+STUDIES?\s*\|\s*(.+)', re.IGNORECASE)
# Detects the "Case Studies" section divider slide (no pipe — not a content slide)
_CASE_STUDI_RE   = re.compile(r'\bcase\s+studi', re.IGNORECASE)
# Detects the Appendix section divider
_APPENDIX_RE     = re.compile(r'^appendix\b', re.IGNORECASE)
# Detects video variant slides
_VIDEO_VARIANT   = re.compile(r'version\s+with\s+video', re.IGNORECASE)
# Matches internal annotation boxes to exclude
_NOTE_LABEL      = re.compile(r'^Note:', re.IGNORECASE)


# ── Slide text helpers ─────────────────────────────────────────────────────────

def _get_slide_texts(slide):
    """Return non-empty text from all shapes, excluding 'Note:' annotation boxes."""
    texts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text and not _NOTE_LABEL.match(text):
            texts.append(text)
    return texts


def _is_section_header_slide(texts):
    """True if this slide is a bare 'Case Studies' section divider (no pipe = not a content slide)."""
    combined = "\n".join(texts)
    return (
        bool(_CASE_STUDI_RE.search(combined))
        and not bool(_HEADING_RE.search(combined))
    )


def _extract_heading_title(texts):
    """Extract project name from the first 'CASE STUDIES | Project Name' shape."""
    for t in texts:
        m = _HEADING_RE.match(t)
        if m:
            return m.group(1).split('\n')[0].strip().upper()
    return ""


# ── Video variant deduplication ────────────────────────────────────────────────

def _dedupe_video_variants(slides):
    """
    Mark base slides with has_video=True when a video variant exists on another slide.
    A video variant slide has is_video_variant=True (contains 'Version with Video' text).
    Video variant slides are dropped from the output.
    If a video variant exists with no matching base, it is kept with has_video=True.
    """
    base_names  = {s["title_hint"].lower().strip() for s in slides if not s["is_video_variant"]}
    video_names = {s["title_hint"].lower().strip() for s in slides if s["is_video_variant"]}

    kept = []
    video_only_seen = set()

    for slide in slides:
        name_lower = slide["title_hint"].lower().strip()

        if not slide["is_video_variant"]:
            kept.append({**slide, "has_video": name_lower in video_names})
        else:
            if name_lower not in base_names and name_lower not in video_only_seen:
                logger.info("Video-only slide %d %r — keeping with has_video=True",
                            slide["slide_num"], slide["title_hint"])
                kept.append({**slide, "has_video": True})
                video_only_seen.add(name_lower)
            else:
                logger.info("Dropping video variant slide %d %r — base version exists",
                            slide["slide_num"], slide["title_hint"])

    return kept


# ── Local metadata inference ───────────────────────────────────────────────────

_BRACKET_RE = re.compile(r'\(([^)]+)\)')

# Ordered: multi-word phrases before single words to avoid partial false matches.
_INDUSTRY_KEYWORDS = [
    ("professional services", "Professional Services"),
    ("public sector",         "Public Sector"),
    ("private equity",        "Private Equity"),
    ("life sciences",         "Healthcare & Life Sciences"),
    ("consumer goods",        "Consumer Goods & Retail"),
    ("petrochemical",         "Chemicals"),
    ("pharmaceutical",        "Healthcare & Life Sciences"),
    ("healthcare",            "Healthcare & Life Sciences"),
    ("chemical",              "Chemicals"),
    ("telecoms",              "Telecommunications IT Media & Electronics"),
    ("telecom",               "Telecommunications IT Media & Electronics"),
    ("aerospace",             "Aerospace & Defense"),
    ("automotive",            "Automotive"),
    ("logistics",             "Travel & Transportation"),
    ("railway",               "Travel & Transportation"),
    ("airline",               "Travel & Transportation"),
    ("banking",               "Financial Services"),
    ("finance",               "Financial Services"),
    ("manufacturing",         "Industrial Goods & Services"),
    ("industrial",            "Industrial Goods & Services"),
    ("energy",                "Energy Utilities & Resources"),
    ("retail",                "Consumer Goods & Retail"),
    ("casino",                "Consumer Goods & Retail"),
    ("gaming",                "Consumer Goods & Retail"),
    ("internal",              "Professional Services"),
]

_AI_TYPE_KEYWORDS = [
    # Multi-word phrases first
    ("computer vision",   "Computer Vision"),
    ("image recognition", "Computer Vision"),
    ("machine learning",  "Machine Learning"),
    ("ml model",          "Machine Learning"),
    ("natural language",  "NLP"),
    ("text analysis",     "NLP"),
    ("data platform",     "Data & Analytics"),
    ("platform build",    "Software & Platform"),
    ("operating model",   "Strategy"),
    # Single words / abbreviations
    ("generative",        "Generative AI"),
    ("genai",             "Generative AI"),
    ("copilot",           "Generative AI"),
    ("llm",               "Generative AI"),
    ("gpt",               "Generative AI"),
    ("rag",               "Generative AI"),
    ("forecasting",       "Machine Learning"),
    ("prediction",        "Machine Learning"),
    ("classification",    "Machine Learning"),
    ("sentiment",         "NLP"),
    ("nlp",               "NLP"),
    ("dashboard",         "Data & Analytics"),
    ("analytics",         "Data & Analytics"),
    ("devops",            "Software & Platform"),
    ("software",          "Software & Platform"),
    ("migration",         "Software & Platform"),
    ("portal",            "Software & Platform"),
    ("roadmap",           "Strategy"),
    ("maturity",          "Strategy"),
    ("strategy",          "Strategy"),
    # Short tokens last — word-boundary matched to avoid false positives
    ("bi",                "Data & Analytics"),
    ("app",               "Software & Platform"),
]


def _match_keywords(text, keyword_map):
    """Return the value for the first keyword found as a whole word in text (case-insensitive)."""
    text_lower = text.lower()
    for keyword, value in keyword_map:
        if re.search(r'\b' + re.escape(keyword) + r'\b', text_lower):
            return value
    return None


def _infer_industry(title, slide_text):
    """Extract bracket text from title and map to industry; fall back to full slide text."""
    m = _BRACKET_RE.search(title)
    if m:
        result = _match_keywords(m.group(1), _INDUSTRY_KEYWORDS)
        if result:
            return result
    return _match_keywords(slide_text, _INDUSTRY_KEYWORDS)


def _infer_ai_type(slide_text):
    """Scan slide text for AI type keywords."""
    return _match_keywords(slide_text, _AI_TYPE_KEYWORDS)


def infer_metadata(title, slide_content):
    """Infer industry_full and ai_type locally — no API calls."""
    return {
        "industry_full": _infer_industry(title, slide_content),
        "ai_type":       _infer_ai_type(slide_content),
    }


# ── PPTX parsing ───────────────────────────────────────────────────────────────

_TEMP_PPTX = os.path.join(os.path.dirname(os.path.abspath(__file__)), "temp_casestudies.pptx")


def parse_pptx(pptx_path):
    """
    Open the PPTX and return a list of slide dicts for every slide between
    the 'Case Studies' section divider and the 'Appendix' section divider.
    Each dict: {slide_num, title_hint, slide_content, is_video_variant}.
    """
    if not pptx_path:
        raise ValueError("PPTX_PATH is not configured in .env")

    _RETRIES = 5
    _RETRY_DELAY = 5
    logger.info("Copying PPTX to local temp file")
    for attempt in range(1, _RETRIES + 1):
        try:
            shutil.copy2(pptx_path, _TEMP_PPTX)
            break
        except PermissionError as exc:
            if attempt < _RETRIES:
                logger.warning(
                    "PPTX locked by OneDrive sync (attempt %d/%d) — retrying in %ds",
                    attempt, _RETRIES, _RETRY_DELAY,
                )
                time.sleep(_RETRY_DELAY)
            else:
                logger.error("PPTX still locked after %d attempts: %s", _RETRIES, pptx_path)
                raise ValueError(
                    "OneDrive is syncing the PPTX and has it locked. "
                    "Wait for the sync icon to clear in the taskbar, then try again. "
                    "To avoid this in future, right-click the file in Explorer and "
                    "choose 'Always keep on this device'."
                ) from exc
        except Exception as exc:
            if attempt < _RETRIES:
                logger.warning(
                    "Could not copy PPTX (attempt %d/%d): %s — retrying in %ds",
                    attempt, _RETRIES, exc, _RETRY_DELAY,
                )
                time.sleep(_RETRY_DELAY)
            else:
                logger.error("Failed to copy PPTX after %d attempts: %s", _RETRIES, exc)
                raise

    try:
        prs = Presentation(_TEMP_PPTX)
        all_slides = list(prs.slides)

        # Find the start: first bare "Case Studies" section divider slide
        start_idx = None
        for i, slide in enumerate(all_slides):
            if _is_section_header_slide(_get_slide_texts(slide)):
                start_idx = i + 1  # begin processing from the slide after the divider
                logger.info("Case Studies section divider found at slide %d", i + 1)
                break

        if start_idx is None:
            logger.warning("No 'Case Studies' section divider found in PPTX — nothing to process")
            return []

        result = []
        for i in range(start_idx, len(all_slides)):
            slide_num = i + 1  # 1-indexed
            texts = _get_slide_texts(all_slides[i])

            # Stop at appendix
            if any(_APPENDIX_RE.match(t) for t in texts):
                logger.info("Appendix slide detected at slide %d — stopping", slide_num)
                break

            is_video_variant = any(_VIDEO_VARIANT.search(t) for t in texts)
            project_name = _extract_heading_title(texts)
            combined = "\n\n".join(texts)

            result.append({
                "slide_num":        slide_num,
                "title_hint":       project_name,
                "slide_content":    combined,
                "is_video_variant": is_video_variant,
            })

        logger.info("PPTX loaded: %d slides in Case Studies range (of %d total)",
                    len(result), len(all_slides))
        return result
    finally:
        try:
            os.remove(_TEMP_PPTX)
        except OSError:
            pass


# ── Main sync entry point ──────────────────────────────────────────────────────

def run_sync():
    """Read the PPTX, process all slides in the Case Studies section, and upsert to the database."""
    pptx_path = os.getenv("PPTX_PATH", "")

    raw_slides = parse_pptx(pptx_path)
    slides     = _dedupe_video_variants(raw_slides)

    added = updated = skipped = unchanged = 0

    for slide in slides:
        project_name = slide["title_hint"]
        if not project_name:
            logger.warning(
                "Slide %d has no 'CASE STUDIES |' heading — skipping", slide["slide_num"]
            )
            skipped += 1
            continue

        content_hash = _hash_content(slide["slide_content"])

        existing = get_case_study_by_hash(content_hash)
        if existing is not None:
            if existing["slide_num"] != slide["slide_num"]:
                update_slide_num(existing["id"], slide["slide_num"])
                logger.info(
                    "Slide %r moved: slide_num %d → %d",
                    project_name, existing["slide_num"], slide["slide_num"],
                )
            else:
                logger.debug(
                    "Slide %d %r unchanged (hash %s) — skipping",
                    slide["slide_num"], project_name, content_hash,
                )
            unchanged += 1
            continue

        meta = infer_metadata(project_name, slide["slide_content"])

        action = upsert_case_study(
            title=project_name,
            slide_num=slide["slide_num"],
            industry_full=meta.get("industry_full"),
            ai_type=meta.get("ai_type"),
            slide_content=slide["slide_content"],
            challenge=None,
            approach=None,
            results=None,
            has_video=1 if slide["has_video"] else 0,
            content_hash=content_hash,
        )
        if action == "added":
            added += 1
        else:
            updated += 1

    log_sync_run(added, updated, skipped, warnings=None)
    total = get_case_study_count()
    logger.info(
        "Sync complete: added=%d updated=%d unchanged=%d skipped=%d total=%d",
        added, updated, unchanged, skipped, total,
    )

    return {
        "added":     added,
        "updated":   updated,
        "unchanged": unchanged,
        "skipped":   skipped,
        "total":     total,
        "warnings":  [],
    }
